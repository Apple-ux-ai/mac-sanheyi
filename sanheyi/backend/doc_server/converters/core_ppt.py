#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT转换器核心逻辑
移植自: PPT和Excel/core/converter_ppt.py
"""

import os
import tempfile
import io
from pathlib import Path
from typing import Dict, Any, Optional
from PIL import Image as PILImage

from .base import BaseConverter

class CorePPTConverter(BaseConverter):
    """PPT转换器核心实现"""
    
    def convert(self, input_path: str, output_path: str, **options) -> Dict[str, Any]:
        """
        执行转换
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            **options: 转换选项
                - quality: 图片质量 (high, medium, low)
        """
        try:
            self.validate_input(input_path)
            self.update_progress(input_path, 0)
            
            ext = Path(input_path).suffix.lower()
            
            result = {'success': False, 'error': 'Unsupported format'}
            
            if ext == '.pdf':
                result = self._convert_pdf_to_ppt(input_path, output_path)
            elif ext in ['.xlsx', '.xls']:
                result = self._convert_excel_to_ppt(input_path, output_path)
            elif ext in ['.pptx', '.ppt']:
                result = self._convert_ppt_to_ppt(input_path, output_path)
            elif ext in ['.docx', '.doc']:
                result = self._convert_word_to_ppt(input_path, output_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff', '.webp']:
                result = self._convert_image_to_ppt(input_path, output_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
            
            if not result.get('success'):
                # 清理可能生成的残余文件
                self.cleanup_on_error(output_path)
            
            return result
        
        except Exception as e:
            self.cleanup_on_error(output_path)
            return {'success': False, 'error': f'转换失败: {str(e)}'}

    def _convert_pdf_to_ppt(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """PDF转PPT"""
        self.update_progress(file_path, 5)
        
        try:
            import fitz  # PyMuPDF
            pymupdf_available = True
        except ImportError:
            pymupdf_available = False
            print("[PDF转PPT] PyMuPDF不可用，将使用pdfplumber（仅文本模式）")
        
        try:
            import pdfplumber
        except ImportError as e:
            if not pymupdf_available:
                return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
        
        self.update_progress(file_path, 15)
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            if pymupdf_available:
                doc = fitz.open(file_path)
                total_pages = len(doc)
                
                if total_pages == 0:
                    doc.close()
                    return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
                
                self.update_progress(file_path, 25)
                
                for page_idx in range(total_pages):
                    page = doc[page_idx]
                    
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    mat = fitz.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    
                    img_data = pix.tobytes("png")
                    
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                        tmp_path = tmp_file.name
                        tmp_file.write(img_data)
                    
                    try:
                        img = PILImage.open(tmp_path)
                        img_width, img_height = img.size
                        img.close()
                        
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height
                        
                        margin = Inches(0.3)
                        available_width = slide_width - 2 * margin
                        available_height = slide_height - 2 * margin
                        
                        width_ratio = available_width / Emu(img_width * 914400 / 96)
                        height_ratio = available_height / Emu(img_height * 914400 / 96)
                        scale = min(width_ratio, height_ratio, 1.0)
                        
                        pic_width = Emu(img_width * 914400 / 96 * scale)
                        pic_height = Emu(img_height * 914400 / 96 * scale)
                        
                        left = (slide_width - pic_width) / 2
                        top = (slide_height - pic_height) / 2
                        
                        slide.shapes.add_picture(tmp_path, int(left), int(top), int(pic_width), int(pic_height))
                    finally:
                        try:
                            os.unlink(tmp_path)
                        except:
                            pass
                    
                    progress = 25 + int(((page_idx + 1) / total_pages) * 60)
                    self.update_progress(file_path, progress)
                
                doc.close()
            else:
                with pdfplumber.open(file_path) as pdf:
                    if len(pdf.pages) == 0:
                        return {'success': False, 'error': 'PDF文件为空，没有可转换的页面'}
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    for page_idx, page in enumerate(pdf.pages, 1):
                        slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(slide_layout)
                        
                        text = page.extract_text()
                        
                        if text:
                            left = Inches(0.3)
                            top = Inches(0.3)
                            width = Inches(9.4)
                            height = Inches(6.9)
                            
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = text_box.text_frame
                            text_frame.word_wrap = True
                            
                            text_length = len(text)
                            if text_length > 2000:
                                font_size = Pt(8)
                            elif text_length > 1000:
                                font_size = Pt(9)
                            elif text_length > 500:
                                font_size = Pt(10)
                            else:
                                font_size = Pt(11)
                            
                            p = text_frame.paragraphs[0]
                            p.text = text
                            p.font.size = font_size
                        
                        progress = 25 + int((page_idx / total_pages) * 60)
                        self.update_progress(file_path, progress)
            
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
            
            self.update_progress(file_path, 90)
            prs.save(output_path)
            self.update_progress(file_path, 100)
            
            return {'success': True, 'output_file': output_path}
            
        except Exception as e:
            return {'success': False, 'error': f'PDF转PPT失败: {str(e)}'}

    def _convert_excel_to_ppt(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """Excel转PPT"""
        ext = Path(file_path).suffix.lower()
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
            
        temp_dir = tempfile.mkdtemp()
        wb_xlsx = None
        wb_xls = None
        
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            if ext == '.xlsx':
                try:
                    from openpyxl import load_workbook
                except ImportError as e:
                    return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
                
                wb_xlsx = load_workbook(file_path, data_only=True)
                sheet_names = wb_xlsx.sheetnames
            else: # .xls
                try:
                    import xlrd
                except ImportError as e:
                    return {'success': False, 'error': f'xlrd库导入失败: {str(e)}'}
                
                wb_xls = xlrd.open_workbook(file_path)
                sheet_names = wb_xls.sheet_names()

            image_paths = {}

            # 第一步：全局预检查是否有图片
            for sheet_name in sheet_names:
                if ext == '.xlsx':
                    ws = wb_xlsx[sheet_name]
                    if hasattr(ws, '_images') and ws._images:
                        return {'success': False, 'error': 'Excel 文件中包含图片，目前暂不支持带图片的转换，请上传纯文本的 Excel 格式。'}
                else:
                    ws = wb_xls.sheet_by_name(sheet_name)
                    # 检查 xlrd 模式下的图片 (xls)
                    if hasattr(wb_xls, 'handle_workbook_objects') and wb_xls.handle_workbook_objects:
                        return {'success': False, 'error': 'Excel 文件中包含图片，目前暂不支持带图片的转换，请上传纯文本的 Excel 格式。'}

                # 检查单元格内容是否包含 WPS 的 DISPIMG 公式 (即单元格内嵌图片)
                # 遍历前 100 行 20 列进行快速检查
                m_row = ws.max_row if ext == '.xlsx' else ws.nrows
                m_col = ws.max_column if ext == '.xlsx' else ws.ncols
                for r in range(min(m_row, 100)):
                    for c in range(min(m_col, 20)):
                        if ext == '.xlsx':
                            val = ws.cell(row=r+1, column=c+1).value
                        else:
                            val = ws.cell_value(r, c)
                        if isinstance(val, str) and val.strip().upper().startswith('=DISPIMG'):
                            return {'success': False, 'error': 'Excel 文件中包含内嵌图片 (DISPIMG)，目前暂不支持带图片的转换，请上传纯文本的 Excel 格式。'}

            for sheet_name in sheet_names:
                if wb_xlsx:
                    ws = wb_xlsx[sheet_name]
                    max_row = ws.max_row
                    max_col = ws.max_column
                else:
                    ws = wb_xls.sheet_by_name(sheet_name)
                    max_row = ws.nrows
                    max_col = ws.ncols
                
                if max_row == 0 and max_col == 0:
                    continue
                
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = sheet_name
                title_para.font.size = Pt(24)
                title_para.font.bold = True

                if max_row > 0 and max_col > 0:
                    max_display_rows = 15
                    max_display_cols = 8
                    
                    actual_rows = min(max_row, max_display_rows)
                    actual_cols = min(max_col, max_display_cols)

                    table = slide.shapes.add_table(
                        actual_rows, actual_cols, 
                        Inches(0.5), Inches(1.2), Inches(9), Inches(5.5)
                    ).table
                    
                    for row_idx in range(actual_rows):
                        for col_idx in range(actual_cols):
                            if wb_xlsx:
                                cell_value = ws.cell(row=row_idx + 1, column=col_idx + 1).value
                            else:
                                cell_value = ws.cell_value(row_idx, col_idx)
                                
                            cell_text = str(cell_value) if cell_value is not None else ''
                            
                            if len(cell_text) > 50:
                                cell_text = cell_text[:47] + '...'
                            
                            table.cell(row_idx, col_idx).text = cell_text
                            
                            cell = table.cell(row_idx, col_idx)
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.size = Pt(10)
                            
                            if row_idx == 0:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run.font.bold = True
                    
                    if max_row > max_display_rows or max_col > max_display_cols:
                        note_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.3))
                        note_frame = note_box.text_frame
                        note_para = note_frame.paragraphs[0]
                        note_para.text = f"注: 仅显示前 {actual_rows} 行 × {actual_cols} 列数据"
                        note_para.font.size = Pt(10)
                        note_para.font.italic = True
            
            if wb_xlsx:
                wb_xlsx.close() # Release file lock
                
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.text = "Excel文件为空或无内容"
                title_para.font.size = Pt(24)
            
            prs.save(output_path)
            return {'success': True, 'output_file': output_path}
            
        except Exception as e:
            return {'success': False, 'error': f'Excel转PPT失败: {str(e)}'}
        finally:
            if wb_xlsx:
                wb_xlsx.close()
            import shutil
            shutil.rmtree(temp_dir, ignore_errors=True)

    def _convert_ppt_to_ppt(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """PPT转PPT (主要是格式清理或另存)"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式（旧格式）不支持，请先转换为PPTX格式'}
            
        try:
            from pptx import Presentation
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
            
        try:
            prs = Presentation(file_path)
            prs.save(output_path)
            return {'success': True, 'output_file': output_path}
        except Exception as e:
            return {'success': False, 'error': f'PPT转换失败: {str(e)}'}

    def _convert_word_to_ppt(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """Word转PPT"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式不支持，请先转换为DOCX格式'}
            
        try:
            from docx import Document
        except ImportError as e:
            return {'success': False, 'error': f'python-docx库导入失败: {str(e)}'}
            
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            return {'success': False, 'error': f'python-pptx库导入失败: {str(e)}'}
            
        try:
            doc = Document(file_path)
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            for para in doc.paragraphs:
                if para.text.strip():
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    title_text = para.text.strip()[:50]
                    title.text = title_text
                    
                    content = slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.word_wrap = True
                    text_frame.text = para.text.strip()
            
            if len(prs.slides) == 0:
                slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(slide_layout)
                
            prs.save(output_path)
            return {'success': True, 'output_file': output_path}
            
        except Exception as e:
            return {'success': False, 'error': f'Word转PPT失败: {str(e)}'}

    def _convert_image_to_ppt(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """图片转PPT"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Emu
        except ImportError as e:
            return {'success': False, 'error': f'依赖库缺失: {str(e)}'}
            
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(slide_layout)
            
            img = PILImage.open(file_path)
            img_width, img_height = img.size
            img.close()
            
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            width_ratio = slide_width / Emu(img_width * 914400 / 96)
            height_ratio = slide_height / Emu(img_height * 914400 / 96)
            scale = min(width_ratio, height_ratio, 1.0)
            
            pic_width = Emu(img_width * 914400 / 96 * scale)
            pic_height = Emu(img_height * 914400 / 96 * scale)
            
            left = (slide_width - pic_width) / 2
            top = (slide_height - pic_height) / 2
            
            slide.shapes.add_picture(file_path, int(left), int(top), int(pic_width), int(pic_height))
            
            prs.save(output_path)
            return {'success': True, 'output_file': output_path}
            
        except Exception as e:
            return {'success': False, 'error': f'图片转PPT失败: {str(e)}'}
