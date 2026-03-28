#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Excel转换器核心逻辑
移植自: PPT和Excel/core/converter_excel.py
"""

import os
from pathlib import Path
from typing import Dict, Any, List, Optional

from .base import BaseConverter

class CoreExcelConverter(BaseConverter):
    """Excel转换器核心实现"""
    
    def convert(self, input_path: str, output_path: str, **options) -> Dict[str, Any]:
        """
        执行转换
        
        Args:
            input_path: 输入文件路径
            output_path: 输出文件路径
            **options: 转换选项
                - mode: 转换模式 (per_page, merge)
        """
        try:
            self.validate_input(input_path)
            self.update_progress(input_path, 0)
            
            ext = Path(input_path).suffix.lower()
            
            result = {'success': False, 'error': 'Unsupported format'}
            
            if ext == '.pdf':
                result = self._convert_pdf_to_excel(input_path, output_path, **options)
            elif ext in ['.xlsx', '.xls']:
                result = self._convert_excel_to_excel(input_path, output_path)
            elif ext in ['.pptx', '.ppt']:
                result = self._convert_ppt_to_excel(input_path, output_path, **options)
            elif ext in ['.docx', '.doc']:
                result = self._convert_word_to_excel(input_path, output_path)
            else:
                return {'success': False, 'error': f'不支持的文件格式: {ext}'}
            
            if not result.get('success'):
                self.cleanup_on_error(output_path)
            
            return result
        
        except Exception as e:
            self.cleanup_on_error(output_path)
            return {'success': False, 'error': f'转换失败: {str(e)}'}

    def _extract_tables_improved(self, page):
        """改进的表格提取方法"""
        tables = []
        try:
            # 策略1: 默认
            default_tables = page.extract_tables()
            if default_tables:
                tables.extend(default_tables)
            
            # 策略2: 严格设置
            table_settings = {
                "vertical_strategy": "lines_strict",
                "horizontal_strategy": "lines_strict",
                "min_words_vertical": 1,
                "min_words_horizontal": 1,
                "snap_tolerance": 3,
                "join_tolerance": 3,
                "edge_min_length": 3,
                "intersection_tolerance": 3
            }
            strict_tables = page.extract_tables(table_settings)
            if strict_tables:
                for table in strict_tables:
                    if table not in tables:
                        tables.append(table)
            
            # 策略3: 宽松设置
            loose_settings = {
                "vertical_strategy": "text",
                "horizontal_strategy": "text",
                "snap_tolerance": 5,
                "join_tolerance": 5
            }
            loose_tables = page.extract_tables(loose_settings)
            if loose_tables:
                for table in loose_tables:
                    if table not in tables and self._is_valid_table_data(table):
                        tables.append(table)
        except Exception as e:
            print(f"[表格提取错误] {e}")
        
        cleaned_tables = []
        for table in tables:
            cleaned_table = self._clean_table_data(table)
            if cleaned_table and len(cleaned_table) > 0:
                cleaned_tables.append(cleaned_table)
        
        return cleaned_tables

    def _is_valid_table_data(self, table):
        if not table or len(table) < 2:
            return False
        non_empty_cells = 0
        total_cells = 0
        for row in table:
            if row:
                for cell in row:
                    total_cells += 1
                    if cell and str(cell).strip():
                        non_empty_cells += 1
        return total_cells > 0 and (non_empty_cells / total_cells) >= 0.3

    def _clean_table_data(self, table):
        if not table: return None
        cleaned_table = []
        for row in table:
            if row:
                cleaned_row = []
                for cell in row:
                    if cell is not None:
                        cell_str = str(cell).strip()
                        cell_str = ' '.join(cell_str.split())
                        cleaned_row.append(cell_str)
                    else:
                        cleaned_row.append("")
                if any(cell for cell in cleaned_row):
                    cleaned_table.append(cleaned_row)
        return cleaned_table if cleaned_table else None

    def _extract_structured_text(self, page):
        try:
            text = page.extract_text()
            if not text: return []
            lines = text.split('\n')
            structured_lines = []
            for line in lines:
                line = line.strip()
                if line:
                    if '\t' in line or '  ' in line:
                        parts = []
                        if '\t' in line:
                            parts = [part.strip() for part in line.split('\t') if part.strip()]
                        else:
                            parts = [part.strip() for part in line.split('  ') if part.strip()]
                        if len(parts) > 1:
                            structured_lines.append(' | '.join(parts))
                        else:
                            structured_lines.append(line)
                    else:
                        structured_lines.append(line)
            return structured_lines
        except Exception:
            return []

    def _convert_pdf_to_excel(self, file_path: str, output_path: str, **options) -> Dict[str, Any]:
        """PDF转Excel"""
        self.update_progress(file_path, 5)
        
        try:
            import pdfplumber
        except ImportError as e:
            return {'success': False, 'error': f'pdfplumber库导入失败: {str(e)}'}
            
        try:
            from openpyxl import Workbook
        except ImportError as e:
            return {'success': False, 'error': f'openpyxl库导入失败: {str(e)}'}
            
        mode = options.get('mode', 'per_page')
        # OCR logic omitted for simplicity and dependency safety
        
        self.update_progress(file_path, 15)
        
        try:
            with pdfplumber.open(file_path) as pdf:
                if mode == 'per_page':
                    wb = Workbook()
                    wb.remove(wb.active)
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    for i, page in enumerate(pdf.pages, 1):
                        ws = wb.create_sheet(title=f"第{i}页")
                        tables = self._extract_tables_improved(page)
                        
                        if tables:
                            current_row = 1
                            for table in tables:
                                for row_idx, row in enumerate(table):
                                    for col_idx, cell in enumerate(row):
                                        if cell is not None:
                                            ws.cell(row=current_row + row_idx, column=col_idx + 1, value=str(cell).strip())
                                current_row += len(table) + 2
                        else:
                            text_data = self._extract_structured_text(page)
                            for row_idx, line in enumerate(text_data, 1):
                                ws.cell(row=row_idx, column=1, value=line)
                        
                        progress = 25 + int((i / total_pages) * 60)
                        self.update_progress(file_path, progress)
                    
                    self.update_progress(file_path, 90)
                    wb.save(output_path)
                    
                else: # merge mode
                    wb = Workbook()
                    ws = wb.active
                    ws.title = "合并内容"
                    
                    total_pages = len(pdf.pages)
                    self.update_progress(file_path, 25)
                    
                    row_offset = 1
                    for page_idx, page in enumerate(pdf.pages, 1):
                        if page_idx > 1:
                            ws.cell(row=row_offset, column=1, value=f"--- 第{page_idx}页 ---")
                            row_offset += 1
                        
                        tables = self._extract_tables_improved(page)
                        
                        if tables:
                            for table in tables:
                                for row in table:
                                    for col_idx, cell in enumerate(row):
                                        if cell is not None:
                                            ws.cell(row=row_offset, column=col_idx + 1, value=str(cell).strip())
                                    row_offset += 1
                                row_offset += 1
                        else:
                            text_data = self._extract_structured_text(page)
                            for line in text_data:
                                ws.cell(row=row_offset, column=1, value=line)
                                row_offset += 1
                        
                        row_offset += 1
                        progress = 25 + int((page_idx / total_pages) * 60)
                        self.update_progress(file_path, progress)
                        
                    self.update_progress(file_path, 90)
                    wb.save(output_path)
            
            self.update_progress(file_path, 100)
            return {'success': True, 'output_file': output_path}
            
        except Exception as e:
            return {'success': False, 'error': f'PDF转Excel失败: {str(e)}'}

    def _convert_excel_to_excel(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """Excel转Excel"""
        try:
            from openpyxl import load_workbook, Workbook
            source_wb = load_workbook(file_path, data_only=True)
            wb = Workbook()
            wb.remove(wb.active)
            
            for sheet_name in source_wb.sheetnames:
                source_ws = source_wb[sheet_name]
                target_ws = wb.create_sheet(title=sheet_name)
                for row in source_ws.iter_rows(values_only=False):
                    for cell in row:
                        if cell.value is not None:
                            target_ws.cell(row=cell.row, column=cell.column, value=cell.value)
            
            wb.save(output_path)
            source_wb.close()
            return {'success': True, 'output_file': output_path}
        except Exception as e:
            return {'success': False, 'error': f'Excel转换失败: {str(e)}'}

    def _convert_ppt_to_excel(self, file_path: str, output_path: str, **options) -> Dict[str, Any]:
        """PPT转Excel"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.ppt':
            return {'success': False, 'error': 'PPT格式不支持，请先转换为PPTX格式'}
            
        try:
            from pptx import Presentation
            from openpyxl import Workbook
        except ImportError as e:
            return {'success': False, 'error': f'依赖库缺失: {str(e)}'}
            
        mode = options.get('mode', 'per_page')
        
        try:
            prs = Presentation(file_path)
            wb = Workbook()
            
            if mode == 'per_page':
                wb.remove(wb.active)
                for i, slide in enumerate(prs.slides, 1):
                    ws = wb.create_sheet(title=f"第{i}页")
                    row = 1
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            ws.cell(row=row, column=1, value=shape.text)
                            row += 1
            else:
                ws = wb.active
                ws.title = "合并内容"
                row = 1
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            ws.cell(row=row, column=1, value=shape.text)
                            row += 1
                    row += 1
            
            wb.save(output_path)
            return {'success': True, 'output_file': output_path}
        except Exception as e:
            return {'success': False, 'error': f'PPT转Excel失败: {str(e)}'}

    def _convert_word_to_excel(self, file_path: str, output_path: str) -> Dict[str, Any]:
        """Word转Excel"""
        input_ext = Path(file_path).suffix.lower()
        if input_ext == '.doc':
            return {'success': False, 'error': 'DOC格式不支持，请先转换为DOCX格式'}
            
        try:
            from docx import Document
            from openpyxl import Workbook
        except ImportError as e:
            return {'success': False, 'error': f'依赖库缺失: {str(e)}'}
            
        try:
            doc = Document(file_path)
            wb = Workbook()
            ws = wb.active
            
            row_idx = 1
            
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    col_idx = 1
                    for cell in row.cells:
                        ws.cell(row=row_idx, column=col_idx, value=cell.text.strip())
                        col_idx += 1
                    row_idx += 1
                row_idx += 2
            
            # 提取文本
            for para in doc.paragraphs:
                if para.text.strip():
                    ws.cell(row=row_idx, column=1, value=para.text.strip())
                    row_idx += 1
            
            wb.save(output_path)
            return {'success': True, 'output_file': output_path}
        except Exception as e:
            return {'success': False, 'error': f'Word转Excel失败: {str(e)}'}
