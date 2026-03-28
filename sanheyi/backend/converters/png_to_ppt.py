from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from .base import BaseConverter
from typing import Dict, Any, Union, List

class PNGToPPTConverter(BaseConverter):
    """PNG到PPT转换器"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['ppt', 'pptx']
    
    def convert(self, input_path: Union[str, List[str]], output_path: str, **options) -> Dict[str, Any]:
        """
        将PNG转换为PPT
        
        Args:
            input_path: PNG文件路径
            output_path: PPT文件路径
            **options: 无特殊参数
        """
        try:
            input_paths = input_path if isinstance(input_path, list) else [input_path]
            for path in input_paths:
                self.validate_input(path)

            prs = Presentation()
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            blank_slide_layout = prs.slide_layouts[6]

            for path in input_paths:
                with Image.open(path) as img:
                    img_width, img_height = img.size

                width_ratio = slide_width / img_width
                height_ratio = slide_height / img_height
                scale_ratio = min(width_ratio, height_ratio)

                scaled_width = int(img_width * scale_ratio)
                scaled_height = int(img_height * scale_ratio)
                left = (slide_width - scaled_width) // 2
                top = (slide_height - scaled_height) // 2

                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(path, left, top, width=scaled_width, height=scaled_height)

            prs.save(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'size': self.get_output_size(output_path)
            }
            
        except Exception as e:
            self.cleanup_on_error(output_path)
            raise Exception(f"PNG to PPT conversion failed: {str(e)}")
